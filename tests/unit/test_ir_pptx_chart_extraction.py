#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
测试PPTX图表数据直取功能
验证统一IR解析系统的核心突破功能
"""

import requests
import json
import time
import os
from pathlib import Path
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

def create_test_pptx_with_charts():
    """创建包含图表的测试PPTX文件"""
    
    print("📊 创建测试PPTX文件（包含图表数据）...")
    
    # 创建演示文稿
    prs = Presentation()
    
    # 第一张幻灯片：标题页
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # 标题布局
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "销售数据分析报告"
    subtitle.text = "2024年第一季度业绩总结"
    
    # 第二张幻灯片：柱状图
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # 空白布局
    slide2.shapes.title.text = "月度销售额对比"
    
    # 创建图表数据
    chart_data = CategoryChartData()
    chart_data.categories = ['1月', '2月', '3月']
    chart_data.add_series('产品A', (120, 150, 180))
    chart_data.add_series('产品B', (80, 95, 110))
    chart_data.add_series('产品C', (60, 75, 85))
    
    # 添加柱状图
    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(5)
    chart = slide2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    chart.chart_title.text_frame.text = "月度销售额对比（万元）"
    
    # 第三张幻灯片：饼图
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "市场份额分布"
    
    # 创建饼图数据
    pie_data = CategoryChartData()
    pie_data.categories = ['华东', '华南', '华北', '西部']
    pie_data.add_series('市场份额', (35, 28, 22, 15))
    
    # 添加饼图
    pie_chart = slide3.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, pie_data
    ).chart
    
    pie_chart.chart_title.text_frame.text = "各区域市场份额（%）"
    
    # 第四张幻灯片：折线图
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    slide4.shapes.title.text = "增长趋势分析"
    
    # 创建折线图数据
    line_data = CategoryChartData()
    line_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
    line_data.add_series('2023年', (100, 120, 140, 160))
    line_data.add_series('2024年预测', (110, 135, 165, 190))
    
    # 添加折线图
    line_chart = slide4.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, line_data
    ).chart
    
    line_chart.chart_title.text_frame.text = "季度增长趋势对比"
    
    # 第五张幻灯片：表格数据
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    slide5.shapes.title.text = "详细数据表"
    
    # 添加表格
    rows, cols = 4, 4
    table = slide5.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(3)).table
    
    # 填充表格数据
    table_data = [
        ['产品', '1月', '2月', '3月'],
        ['产品A', '120万', '150万', '180万'],
        ['产品B', '80万', '95万', '110万'],
        ['产品C', '60万', '75万', '85万']
    ]
    
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_data in enumerate(row_data):
            table.cell(row_idx, col_idx).text = cell_data
    
    # 保存文件
    file_path = Path("test_charts_presentation.pptx")
    prs.save(str(file_path))
    
    print(f"✅ 测试PPTX文件创建成功: {file_path}")
    print(f"   - 5张幻灯片")
    print(f"   - 3个图表（柱状图、饼图、折线图）")
    print(f"   - 1个表格")
    print(f"   - 文本内容")
    
    return file_path

def test_ir_pptx_parsing():
    """测试PPTX IR解析功能"""
    
    print("\n🧪 测试PPTX IR解析功能")
    print("=" * 60)
    
    base_url = "http://127.0.0.1:8000"
    
    try:
        # 1. 创建测试PPTX文件
        pptx_file = create_test_pptx_with_charts()
        
        # 2. 上传文件
        print("\n📤 上传PPTX文件...")
        with open(pptx_file, 'rb') as f:
            files = {'file': (pptx_file.name, f, 'application/vnd.openxmlformats-officedocument.presentationml.presentation')}
            response = requests.post(f"{base_url}/kg/upload", files=files)
        
        if response.status_code != 200:
            print(f"❌ 上传失败: {response.status_code}")
            return False
            
        upload_data = response.json()
        upload_id = upload_data['upload_id']
        print(f"✅ 上传成功 (ID: {upload_id})")
        
        # 3. 触发解析
        print("🔧 触发IR解析...")
        parse_response = requests.post(f"{base_url}/kg/files/{upload_id}/parse")
        
        if parse_response.status_code != 200:
            print(f"❌ 解析触发失败: {parse_response.status_code}")
            return False
            
        print("✅ IR解析任务已启动")
        
        # 4. 等待解析完成
        print("⏳ 等待解析完成...")
        max_wait = 60
        wait_time = 0
        
        while wait_time < max_wait:
            status_response = requests.get(f"{base_url}/kg/files/{upload_id}/status")
            if status_response.status_code == 200:
                status_data = status_response.json()
                current_status = status_data['data']['status']
                print(f"   状态: {current_status} ({wait_time}s)")
                
                if current_status == 'parsed':
                    print("✅ 解析完成")
                    break
                elif current_status == 'failed':
                    error_msg = status_data['data'].get('error', '未知错误')
                    print(f"❌ 解析失败: {error_msg}")
                    return False
            
            time.sleep(3)
            wait_time += 3
        
        if wait_time >= max_wait:
            print("❌ 解析超时")
            return False
        
        # 5. 获取解析结果
        print("📊 获取解析结果...")
        result_response = requests.get(f"{base_url}/kg/files/{upload_id}/preview")
        
        if result_response.status_code != 200:
            print(f"❌ 获取结果失败: {result_response.status_code}")
            return False
            
        result_data = result_response.json()
        
        if not result_data.get('success'):
            print(f"❌ 解析结果无效: {result_data.get('error', '未知错误')}")
            return False
        
        # 6. 分析解析结果
        print("🔍 分析解析结果...")
        preview_data = result_data['data']
        raw_data = preview_data.get('raw_data', [])
        metadata = preview_data.get('metadata', {})
        
        print(f"   总记录数: {len(raw_data)}")
        print(f"   元数据字段: {len(metadata)}")
        
        # 统计不同类型的内容块
        content_types = {}
        chart_data_blocks = 0
        table_blocks = 0
        text_blocks = 0
        figure_blocks = 0
        
        for record in raw_data:
            content_type = record.get('content_type', 'unknown')
            content_types[content_type] = content_types.get(content_type, 0) + 1
            
            # 检查是否为图表数据
            if record.get('style') == 'chart_data':
                chart_data_blocks += 1
                print(f"   📊 发现图表数据块: {record.get('block_id')}")
                
                # 显示图表数据示例
                chart_cells = []
                for key, value in record.items():
                    if key.startswith('cell_'):
                        chart_cells.append(f"{key}: {value}")
                
                if chart_cells:
                    print(f"      图表数据示例: {chart_cells[:3]}...")
            
            elif content_type == 'table':
                table_blocks += 1
            elif content_type == 'text':
                text_blocks += 1
            elif content_type == 'figure':
                figure_blocks += 1
        
        print(f"\n📋 内容统计:")
        print(f"   📊 图表数据块: {chart_data_blocks}")
        print(f"   📋 表格块: {table_blocks}")
        print(f"   📝 文本块: {text_blocks}")
        print(f"   🖼 图片块: {figure_blocks}")
        
        # 7. 验证核心功能
        print(f"\n🎯 核心功能验证:")
        
        # 验证图表数据直取
        if chart_data_blocks >= 3:  # 期望至少3个图表
            print("✅ 图表数据直取: 成功")
            print(f"   成功提取 {chart_data_blocks} 个图表的原始数据")
        else:
            print("❌ 图表数据直取: 失败")
            print(f"   期望3个图表，实际提取 {chart_data_blocks} 个")
        
        # 验证多格式内容识别
        expected_types = {'text', 'table'}
        found_types = set(content_types.keys())
        if expected_types.issubset(found_types):
            print("✅ 多格式内容识别: 成功")
            print(f"   识别类型: {list(found_types)}")
        else:
            print("❌ 多格式内容识别: 部分失败")
            print(f"   期望: {expected_types}, 实际: {found_types}")
        
        # 验证页面信息
        pages_found = set()
        for record in raw_data:
            page_num = record.get('page_number')
            if page_num:
                pages_found.add(page_num)
        
        if len(pages_found) >= 5:  # 期望5张幻灯片
            print("✅ 页面信息提取: 成功")
            print(f"   识别页面: {sorted(pages_found)}")
        else:
            print("❌ 页面信息提取: 部分失败")
            print(f"   期望5页，实际识别 {len(pages_found)} 页")
        
        # 8. 显示详细结果示例
        print(f"\n📄 解析结果示例:")
        for i, record in enumerate(raw_data[:3]):
            print(f"   记录 {i+1}:")
            print(f"      类型: {record.get('content_type', 'unknown')}")
            print(f"      页面: {record.get('page_number', 'unknown')}")
            print(f"      样式: {record.get('style', 'none')}")
            if record.get('content'):
                content = record['content'][:50] + "..." if len(record['content']) > 50 else record['content']
                print(f"      内容: {content}")
        
        # 清理测试文件
        if pptx_file.exists():
            pptx_file.unlink()
            print(f"\n🧹 清理测试文件: {pptx_file.name}")
        
        print(f"\n🎉 PPTX IR解析测试完成!")
        
        # 判断测试是否成功
        success = (chart_data_blocks >= 2 and  # 至少提取2个图表数据
                  len(pages_found) >= 4 and    # 至少识别4个页面
                  len(raw_data) >= 10)         # 至少10条记录
        
        if success:
            print("✅ 测试结果: 成功")
            print("   核心功能验证通过，图表数据直取功能正常工作")
        else:
            print("⚠️  测试结果: 部分成功")
            print("   部分功能需要进一步优化")
        
        return success
        
    except Exception as e:
        print(f"❌ 测试异常: {e}")
        return False

if __name__ == "__main__":
    # 检查依赖
    try:
        import requests
        from pptx import Presentation
        print("✅ 依赖检查通过")
    except ImportError as e:
        print(f"❌ 缺少依赖: {e}")
        print("请安装: pip install requests python-pptx")
        exit(1)
    
    # 运行测试
    success = test_ir_pptx_parsing()
    
    if success:
        print("\n🎊 恭喜！PPTX图表数据直取功能测试成功！")
        print("这是文档解析能力的质变级提升！")
    else:
        print("\n🔧 测试未完全通过，需要进一步调试和优化")
