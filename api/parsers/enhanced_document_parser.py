#!/usr/bin/env python3
"""
增强的多格式文档解析器
支持PDF、Word、PowerPoint、文本等多种格式
"""

import logging
import re
from pathlib import Path
from typing import List, Dict, Any, Optional
import pandas as pd

logger = logging.getLogger(__name__)

class EnhancedDocumentParser:
    """增强的文档解析器"""
    
    def __init__(self):
        self.supported_formats = {
            '.pdf': self.parse_pdf_enhanced,
            '.docx': self.parse_docx_enhanced,
            '.doc': self.parse_doc_enhanced,
            '.pptx': self.parse_pptx_enhanced,
            '.ppt': self.parse_ppt_enhanced,
            '.txt': self.parse_text_enhanced,
            '.csv': self.parse_csv_enhanced,
            '.md': self.parse_markdown_enhanced,
            '.rtf': self.parse_rtf_enhanced
        }
    
    def parse_document(self, file_path: Path) -> Dict[str, Any]:
        """
        解析文档，返回结构化数据
        
        Args:
            file_path: 文档文件路径
            
        Returns:
            解析结果字典
        """
        file_ext = file_path.suffix.lower()
        
        if file_ext not in self.supported_formats:
            raise ValueError(f"不支持的文件格式: {file_ext}")
        
        logger.info(f"开始解析文档: {file_path} (格式: {file_ext})")
        
        try:
            parser_func = self.supported_formats[file_ext]
            result = parser_func(file_path)
            
            # 统一返回格式
            return {
                'success': True,
                'data': {
                    'raw_data': result.get('raw_data', []),
                    'entities': result.get('entities', []),
                    'relations': result.get('relations', []),
                    'metadata': result.get('metadata', {})
                },
                'stats': {
                    'total_records': len(result.get('raw_data', [])),
                    'total_entities': len(result.get('entities', [])),
                    'total_relations': len(result.get('relations', [])),
                    'file_type': file_ext[1:],
                    'file_size': file_path.stat().st_size if file_path.exists() else 0
                }
            }
            
        except Exception as e:
            logger.error(f"文档解析失败: {e}")
            return {
                'success': False,
                'error': str(e),
                'data': None
            }
    
    def parse_pdf_enhanced(self, file_path: Path) -> Dict[str, Any]:
        """增强的PDF解析"""
        try:
            import pdfplumber
            import PyPDF2
            
            raw_data = []
            entities = []
            metadata = {}
            
            # 使用pdfplumber提取文本和表格
            with pdfplumber.open(str(file_path)) as pdf:
                metadata['total_pages'] = len(pdf.pages)
                
                for page_num, page in enumerate(pdf.pages, 1):
                    # 提取文本
                    text = page.extract_text()
                    if text:
                        # 清理文本
                        cleaned_text = self._clean_pdf_text(text)
                        if cleaned_text:
                            paragraphs = self._split_into_paragraphs(cleaned_text)
                            for para_num, para in enumerate(paragraphs, 1):
                                if para.strip() and len(para.strip()) > 3:  # 过滤太短的段落
                                    record = {
                                        '_row_number': len(raw_data) + 1,
                                        'page_number': page_num,
                                        'paragraph_number': para_num,
                                        'content_type': 'text',
                                        'content': para.strip(),
                                        'word_count': len(para.split()),
                                        'char_count': len(para)
                                    }
                                    raw_data.append(record)
                    
                    # 提取表格
                    tables = page.extract_tables()
                    for table_num, table in enumerate(tables, 1):
                        if table:
                            # 处理表格数据
                            headers = table[0] if table else []
                            for row_num, row in enumerate(table[1:], 1):
                                if row and any(cell for cell in row if cell):
                                    record = {
                                        '_row_number': len(raw_data) + 1,
                                        'page_number': page_num,
                                        'table_number': table_num,
                                        'row_number': row_num,
                                        'content_type': 'table',
                                        'content': ' | '.join(str(cell) if cell else '' for cell in row)
                                    }
                                    # 添加表格列数据
                                    for col_num, (header, cell) in enumerate(zip(headers, row), 1):
                                        if header and cell:
                                            record[f'column_{col_num}_{header}'] = str(cell)
                                    
                                    raw_data.append(record)
            
            # 提取实体
            entities = self._extract_entities_from_text_data(raw_data)
            
            return {
                'raw_data': raw_data,
                'entities': entities,
                'relations': [],
                'metadata': metadata
            }
            
        except ImportError:
            logger.error("PDF解析依赖未安装，请安装: pip install pdfplumber PyPDF2")
            raise Exception("PDF解析依赖未安装")
        except Exception as e:
            logger.error(f"PDF解析失败: {e}")
            raise
    
    def parse_docx_enhanced(self, file_path: Path) -> Dict[str, Any]:
        """增强的Word文档解析"""
        try:
            import docx
            
            raw_data = []
            entities = []
            metadata = {}
            
            doc = docx.Document(str(file_path))
            
            # 提取段落
            for para_num, para in enumerate(doc.paragraphs, 1):
                if para.text.strip():
                    record = {
                        '_row_number': len(raw_data) + 1,
                        'paragraph_number': para_num,
                        'content_type': 'paragraph',
                        'content': para.text.strip(),
                        'style': para.style.name if para.style else 'Normal',
                        'word_count': len(para.text.split()),
                        'char_count': len(para.text)
                    }
                    raw_data.append(record)
            
            # 提取表格
            for table_num, table in enumerate(doc.tables, 1):
                headers = []
                if table.rows:
                    headers = [cell.text.strip() for cell in table.rows[0].cells]
                
                for row_num, row in enumerate(table.rows[1:], 1):
                    row_data = [cell.text.strip() for cell in row.cells]
                    if any(cell for cell in row_data if cell):
                        record = {
                            '_row_number': len(raw_data) + 1,
                            'table_number': table_num,
                            'row_number': row_num,
                            'content_type': 'table',
                            'content': ' | '.join(row_data)
                        }
                        # 添加表格列数据
                        for col_num, (header, cell) in enumerate(zip(headers, row_data), 1):
                            if header and cell:
                                record[f'column_{col_num}_{header}'] = cell
                        
                        raw_data.append(record)
            
            metadata['total_paragraphs'] = len([p for p in doc.paragraphs if p.text.strip()])
            metadata['total_tables'] = len(doc.tables)
            
            # 提取实体
            entities = self._extract_entities_from_text_data(raw_data)
            
            return {
                'raw_data': raw_data,
                'entities': entities,
                'relations': [],
                'metadata': metadata
            }
            
        except ImportError:
            logger.error("Word解析依赖未安装，请安装: pip install python-docx")
            raise Exception("Word解析依赖未安装")
        except Exception as e:
            logger.error(f"Word文档解析失败: {e}")
            raise
    
    def parse_pptx_enhanced(self, file_path: Path) -> Dict[str, Any]:
        """增强的PowerPoint解析"""
        try:
            import pptx
            
            raw_data = []
            entities = []
            metadata = {}
            
            ppt = pptx.Presentation(str(file_path))
            metadata['total_slides'] = len(ppt.slides)
            
            for slide_num, slide in enumerate(ppt.slides, 1):
                # 提取文本框内容
                for shape_num, shape in enumerate(slide.shapes, 1):
                    if hasattr(shape, 'text') and shape.text.strip():
                        record = {
                            '_row_number': len(raw_data) + 1,
                            'slide_number': slide_num,
                            'shape_number': shape_num,
                            'content_type': 'text',
                            'content': shape.text.strip(),
                            'word_count': len(shape.text.split()),
                            'char_count': len(shape.text)
                        }
                        raw_data.append(record)
                    
                    # 提取表格
                    if shape.shape_type == 19:  # Table
                        table = shape.table
                        headers = [cell.text.strip() for cell in table.rows[0].cells]
                        
                        for row_num, row in enumerate(table.rows[1:], 1):
                            row_data = [cell.text.strip() for cell in row.cells]
                            if any(cell for cell in row_data if cell):
                                record = {
                                    '_row_number': len(raw_data) + 1,
                                    'slide_number': slide_num,
                                    'table_number': shape_num,
                                    'row_number': row_num,
                                    'content_type': 'table',
                                    'content': ' | '.join(row_data)
                                }
                                # 添加表格列数据
                                for col_num, (header, cell) in enumerate(zip(headers, row_data), 1):
                                    if header and cell:
                                        record[f'column_{col_num}_{header}'] = cell
                                
                                raw_data.append(record)
            
            # 提取实体
            entities = self._extract_entities_from_text_data(raw_data)
            
            return {
                'raw_data': raw_data,
                'entities': entities,
                'relations': [],
                'metadata': metadata
            }
            
        except ImportError:
            logger.error("PowerPoint解析依赖未安装，请安装: pip install python-pptx")
            raise Exception("PowerPoint解析依赖未安装")
        except Exception as e:
            logger.error(f"PowerPoint解析失败: {e}")
            raise
    
    def parse_csv_enhanced(self, file_path: Path) -> Dict[str, Any]:
        """增强的CSV解析"""
        try:
            # 尝试不同编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16']
            df = None
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if df is None:
                raise Exception("无法读取CSV文件，编码不支持")
            
            raw_data = []
            for i, row in df.iterrows():
                record = {'_row_number': i + 2}  # CSV行号从2开始（第1行是标题）
                
                for col_name in df.columns:
                    value = row[col_name]
                    if pd.isna(value):
                        record[col_name] = None
                    else:
                        record[col_name] = str(value).strip()
                
                raw_data.append(record)
            
            # 提取实体
            entities = self._extract_entities_from_structured_data(raw_data)
            
            return {
                'raw_data': raw_data,
                'entities': entities,
                'relations': [],
                'metadata': {
                    'total_rows': len(df),
                    'total_columns': len(df.columns),
                    'columns': list(df.columns)
                }
            }
            
        except Exception as e:
            logger.error(f"CSV解析失败: {e}")
            raise
    
    def parse_text_enhanced(self, file_path: Path) -> Dict[str, Any]:
        """增强的文本文件解析"""
        try:
            # 尝试不同编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16']
            content = None
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            if content is None:
                raise Exception("无法读取文本文件，编码不支持")
            
            raw_data = []
            paragraphs = self._split_into_paragraphs(content)
            
            for para_num, para in enumerate(paragraphs, 1):
                if para.strip():
                    record = {
                        '_row_number': para_num,
                        'paragraph_number': para_num,
                        'content_type': 'text',
                        'content': para.strip(),
                        'word_count': len(para.split()),
                        'char_count': len(para)
                    }
                    raw_data.append(record)
            
            # 提取实体
            entities = self._extract_entities_from_text_data(raw_data)
            
            return {
                'raw_data': raw_data,
                'entities': entities,
                'relations': [],
                'metadata': {
                    'total_paragraphs': len(paragraphs),
                    'total_chars': len(content),
                    'total_words': len(content.split())
                }
            }
            
        except Exception as e:
            logger.error(f"文本文件解析失败: {e}")
            raise
    
    def parse_markdown_enhanced(self, file_path: Path) -> Dict[str, Any]:
        """增强的Markdown解析"""
        return self.parse_text_enhanced(file_path)  # 暂时使用文本解析
    
    def parse_rtf_enhanced(self, file_path: Path) -> Dict[str, Any]:
        """增强的RTF解析"""
        return self.parse_text_enhanced(file_path)  # 暂时使用文本解析
    
    def parse_doc_enhanced(self, file_path: Path) -> Dict[str, Any]:
        """增强的旧版Word解析"""
        # 旧版Word文档需要特殊处理，暂时返回基础解析
        return self.parse_text_enhanced(file_path)
    
    def parse_ppt_enhanced(self, file_path: Path) -> Dict[str, Any]:
        """增强的旧版PowerPoint解析"""
        # 旧版PPT需要特殊处理，暂时返回基础解析
        return self.parse_text_enhanced(file_path)
    
    def _clean_pdf_text(self, text: str) -> str:
        """清理PDF提取的文本"""
        if not text:
            return ""

        # 移除过多的换行符
        text = re.sub(r'\n{3,}', '\n\n', text)

        # 移除单独的字符行（通常是乱码）
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            line = line.strip()
            # 过滤掉只有单个字符或者全是特殊字符的行
            if len(line) > 2 and not re.match(r'^[^\w\s]*$', line):
                cleaned_lines.append(line)

        return '\n'.join(cleaned_lines)

    def _split_into_paragraphs(self, text: str) -> List[str]:
        """将文本分割为段落"""
        # 按双换行符分割段落
        paragraphs = re.split(r'\n\s*\n', text)
        return [para.strip() for para in paragraphs if para.strip()]
    
    def _extract_entities_from_text_data(self, raw_data: List[Dict]) -> List[Dict]:
        """从文本数据中提取实体"""
        entities = []
        
        for i, record in enumerate(raw_data):
            content = record.get('content', '')
            if content:
                # 简单的实体提取
                entities.append({
                    'id': f"text_entity_{i}",
                    'type': 'text_content',
                    'value': content[:100],  # 截取前100字符
                    'source_record': i,
                    'confidence': 1.0
                })
        
        return entities
    
    def _extract_entities_from_structured_data(self, raw_data: List[Dict]) -> List[Dict]:
        """从结构化数据中提取实体"""
        entities = []
        
        for i, record in enumerate(raw_data):
            for field, value in record.items():
                if field.startswith('_'):  # 跳过系统字段
                    continue
                if value and str(value).strip():
                    entities.append({
                        'id': f"field_{field}_{i}",
                        'type': field,
                        'value': str(value),
                        'source_record': i,
                        'confidence': 1.0
                    })
        
        return entities


# 兼容性函数
def parse_document_enhanced(file_path: Path) -> Dict[str, Any]:
    """
    增强的文档解析函数
    """
    parser = EnhancedDocumentParser()
    return parser.parse_document(file_path)
