#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
增强PPTX解析器
支持图表数据直取 + OCR兜底
这是质变级功能提升的核心模块
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import hashlib
import logging
from .ir_core import DocumentIR, IRBlock, BlockType, create_paragraph_block, create_table_block, create_figure_block
from .ocr_engine import get_ocr_engine

logger = logging.getLogger(__name__)

class EnhancedPPTXParser:
    """增强PPTX解析器"""
    
    def __init__(self, use_ocr: bool = True, ocr_confidence: float = 0.5):
        """
        初始化解析器
        
        Args:
            use_ocr: 是否启用OCR
            ocr_confidence: OCR置信度阈值
        """
        self.use_ocr = use_ocr
        self.ocr_confidence = ocr_confidence
        self.ocr_engine = get_ocr_engine() if use_ocr else None
    
    def parse(self, file_path: Path, output_dir: Path) -> DocumentIR:
        """
        解析PPTX文件为统一IR格式
        
        Args:
            file_path: PPTX文件路径
            output_dir: 输出目录（用于保存提取的图片）
            
        Returns:
            DocumentIR: 统一中间表示
        """
        logger.info(f"开始解析PPTX文件: {file_path}")
        
        # 创建输出目录
        output_dir.mkdir(parents=True, exist_ok=True)
        
        try:
            # 加载演示文稿
            prs = Presentation(str(file_path))
            
            # 构建元信息
            meta = {
                "file_id": file_path.stem,
                "type": "pptx",
                "pages": len(prs.slides),
                "title": self._extract_title(prs),
                "author": self._extract_author(prs),
                "created": self._extract_created_time(prs)
            }
            
            logger.info(f"PPTX基本信息: {len(prs.slides)}张幻灯片")
            
            # 解析所有幻灯片
            blocks = []
            block_id_counter = 0
            
            for slide_idx, slide in enumerate(prs.slides, start=1):
                logger.debug(f"解析第{slide_idx}张幻灯片")
                
                # 处理文本框
                text_blocks, block_id_counter = self._extract_text_blocks(
                    slide, slide_idx, block_id_counter
                )
                blocks.extend(text_blocks)
                
                # 处理表格
                table_blocks, block_id_counter = self._extract_table_blocks(
                    slide, slide_idx, block_id_counter
                )
                blocks.extend(table_blocks)
                
                # 处理图表（核心突破功能）
                chart_blocks, block_id_counter = self._extract_chart_blocks(
                    slide, slide_idx, block_id_counter
                )
                blocks.extend(chart_blocks)
                
                # 处理图片
                image_blocks, block_id_counter = self._extract_image_blocks(
                    slide, slide_idx, block_id_counter, output_dir
                )
                blocks.extend(image_blocks)
            
            logger.info(f"PPTX解析完成: 共{len(blocks)}个内容块")
            
            return DocumentIR(meta=meta, blocks=blocks)
            
        except Exception as e:
            logger.error(f"PPTX解析失败: {e}")
            raise
    
    def _extract_title(self, prs: Presentation) -> str:
        """提取演示文稿标题"""
        try:
            if prs.core_properties.title:
                return prs.core_properties.title
            
            # 尝试从第一张幻灯片的标题获取
            if prs.slides:
                first_slide = prs.slides[0]
                for shape in first_slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        return shape.text_frame.text.strip()
            
            return "未知标题"
        except:
            return "未知标题"
    
    def _extract_author(self, prs: Presentation) -> str:
        """提取作者信息"""
        try:
            return prs.core_properties.author or "未知作者"
        except:
            return "未知作者"
    
    def _extract_created_time(self, prs: Presentation) -> str:
        """提取创建时间"""
        try:
            created = prs.core_properties.created
            return created.isoformat() if created else ""
        except:
            return ""
    
    def _extract_text_blocks(self, slide, page: int, start_id: int) -> Tuple[List[IRBlock], int]:
        """提取文本框内容"""
        blocks = []
        block_id = start_id
        
        for shape in slide.shapes:
            if self._is_text_shape(shape):
                try:
                    # 提取文本内容
                    text_content = self._extract_shape_text(shape)
                    if text_content.strip():
                        block_id += 1
                        
                        # 获取样式信息
                        style_info = self._extract_text_style(shape)
                        
                        block = create_paragraph_block(
                            block_id=f"text_{block_id}",
                            page=page,
                            text=text_content,
                            style=style_info
                        )
                        blocks.append(block)
                        
                        logger.debug(f"提取文本块: {len(text_content)}字符")
                        
                except Exception as e:
                    logger.warning(f"文本提取失败: {e}")
                    continue
        
        return blocks, block_id
    
    def _extract_table_blocks(self, slide, page: int, start_id: int) -> Tuple[List[IRBlock], int]:
        """提取表格内容"""
        blocks = []
        block_id = start_id
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                try:
                    # 提取表格数据
                    table_data = self._extract_table_data(shape.table)
                    if table_data:
                        block_id += 1
                        
                        block = create_table_block(
                            block_id=f"table_{block_id}",
                            page=page,
                            cells=table_data,
                            style="pptx_table"
                        )
                        blocks.append(block)
                        
                        logger.debug(f"提取表格: {len(table_data)}行 x {len(table_data[0]) if table_data else 0}列")
                        
                except Exception as e:
                    logger.warning(f"表格提取失败: {e}")
                    continue
        
        return blocks, block_id
    
    def _extract_chart_blocks(self, slide, page: int, start_id: int) -> Tuple[List[IRBlock], int]:
        """
        提取图表数据（核心突破功能）
        直接从图表对象获取原始数据，无需OCR
        """
        blocks = []
        block_id = start_id
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART and hasattr(shape, 'chart'):
                try:
                    chart = shape.chart
                    logger.debug(f"发现图表对象: {chart.chart_type}")
                    
                    # 提取图表数据
                    chart_data = self._extract_chart_data(chart)
                    
                    if chart_data and len(chart_data) > 1:  # 至少有标题行和数据行
                        block_id += 1
                        
                        block = create_table_block(
                            block_id=f"chart_{block_id}",
                            page=page,
                            cells=chart_data,
                            style="chart_data"
                        )
                        
                        # 添加图表特有的元数据
                        block.metadata = {
                            "chart_type": str(chart.chart_type),
                            "has_legend": hasattr(chart, 'has_legend') and chart.has_legend,
                            "series_count": len(chart.series) if chart.series else 0
                        }
                        
                        blocks.append(block)
                        
                        logger.info(f"成功提取图表数据: {len(chart_data)}行 x {len(chart_data[0])}列")
                        
                except Exception as e:
                    logger.warning(f"图表数据提取失败: {e}")
                    # 图表数据提取失败时，将其作为图片处理
                    continue
        
        return blocks, block_id
    
    def _extract_image_blocks(self, slide, page: int, start_id: int, output_dir: Path) -> Tuple[List[IRBlock], int]:
        """提取图片并进行OCR识别"""
        blocks = []
        block_id = start_id
        
        for shape in slide.shapes:
            if self._is_image_shape(shape):
                try:
                    # 保存图片
                    image_path = self._save_shape_image(shape, page, output_dir)
                    if not image_path:
                        continue
                    
                    # OCR识别
                    ocr_text = ""
                    confidence = 0.0
                    
                    if self.use_ocr and self.ocr_engine:
                        ocr_text = self.ocr_engine.extract_text(str(image_path), self.ocr_confidence)
                        confidence = self.ocr_engine.get_text_confidence(str(image_path))
                    
                    # 判断图片类型
                    figure_type = self._classify_figure_type(shape, ocr_text)
                    
                    block_id += 1
                    
                    block = create_figure_block(
                        block_id=f"figure_{block_id}",
                        page=page,
                        image_path=str(image_path),
                        ocr_text=ocr_text,
                        figure_type=figure_type,
                        confidence=confidence
                    )
                    
                    blocks.append(block)
                    
                    logger.debug(f"提取图片: {image_path.name}, OCR: {len(ocr_text)}字符")
                    
                except Exception as e:
                    logger.warning(f"图片处理失败: {e}")
                    continue
        
        return blocks, block_id
    
    def _extract_chart_data(self, chart) -> List[List[str]]:
        """
        从图表对象提取原始数据
        这是PPTX解析的核心突破功能
        """
        try:
            chart_data = []
            
            # 获取分类轴标签
            categories = []
            if hasattr(chart, 'category_axis') and chart.category_axis:
                try:
                    # 尝试获取分类名称
                    if hasattr(chart.category_axis, 'category_names'):
                        categories = [str(cat) for cat in chart.category_axis.category_names if cat]
                except:
                    pass
            
            # 如果没有分类轴，尝试从系列中获取
            if not categories and chart.series:
                try:
                    first_series = chart.series[0]
                    if hasattr(first_series, 'categories'):
                        categories = [str(cat) for cat in first_series.categories if cat is not None]
                except:
                    pass
            
            # 获取系列数据
            series_data = []
            series_names = []
            
            for series in chart.series:
                try:
                    series_name = series.name if series.name else f"系列{len(series_names) + 1}"
                    series_names.append(str(series_name))
                    
                    # 获取数值
                    values = []
                    if hasattr(series, 'values'):
                        for val in series.values:
                            if val is not None:
                                values.append(str(val))
                            else:
                                values.append("")
                    
                    series_data.append(values)
                    
                except Exception as e:
                    logger.warning(f"系列数据提取失败: {e}")
                    continue
            
            # 构建表格数据
            if series_data:
                # 创建表头
                header = ["类别"] + series_names
                chart_data.append(header)
                
                # 创建数据行
                max_rows = max(len(categories), max(len(series) for series in series_data) if series_data else 0)
                
                for i in range(max_rows):
                    row = []
                    
                    # 添加分类
                    if i < len(categories):
                        row.append(categories[i])
                    else:
                        row.append(f"项目{i+1}")
                    
                    # 添加各系列的值
                    for series in series_data:
                        if i < len(series):
                            row.append(series[i])
                        else:
                            row.append("")
                    
                    chart_data.append(row)
            
            return chart_data
            
        except Exception as e:
            logger.error(f"图表数据提取异常: {e}")
            return []
    
    def _is_text_shape(self, shape) -> bool:
        """判断是否为文本形状"""
        return (hasattr(shape, 'has_text_frame') and 
                shape.has_text_frame and 
                shape.shape_type != MSO_SHAPE_TYPE.TABLE and
                shape.shape_type != MSO_SHAPE_TYPE.CHART)
    
    def _is_image_shape(self, shape) -> bool:
        """判断是否为图片形状"""
        return (hasattr(shape, 'image') or 
                shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
    
    def _extract_shape_text(self, shape) -> str:
        """提取形状中的文本"""
        try:
            if hasattr(shape, 'text_frame'):
                paragraphs = []
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        paragraphs.append(paragraph.text.strip())
                return "\n".join(paragraphs)
            elif hasattr(shape, 'text'):
                return shape.text.strip()
            return ""
        except:
            return ""
    
    def _extract_text_style(self, shape) -> str:
        """提取文本样式信息"""
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                first_para = shape.text_frame.paragraphs[0]
                if hasattr(first_para, 'font'):
                    font = first_para.font
                    style_info = []
                    if hasattr(font, 'name') and font.name:
                        style_info.append(f"font:{font.name}")
                    if hasattr(font, 'size') and font.size:
                        style_info.append(f"size:{font.size}")
                    if hasattr(font, 'bold') and font.bold:
                        style_info.append("bold")
                    return ";".join(style_info)
            return "normal"
        except:
            return "normal"
    
    def _extract_table_data(self, table) -> List[List[str]]:
        """提取表格数据"""
        try:
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text_frame.text.strip() if cell.text_frame else ""
                    row_data.append(cell_text)
                table_data.append(row_data)
            return table_data
        except Exception as e:
            logger.warning(f"表格数据提取失败: {e}")
            return []
    
    def _save_shape_image(self, shape, page: int, output_dir: Path) -> Optional[Path]:
        """保存形状中的图片"""
        try:
            if hasattr(shape, 'image'):
                image_blob = shape.image.blob
                image_hash = hashlib.md5(image_blob).hexdigest()[:8]
                image_path = output_dir / f"slide_{page}_{image_hash}.png"
                image_path.write_bytes(image_blob)
                return image_path
            return None
        except Exception as e:
            logger.warning(f"图片保存失败: {e}")
            return None
    
    def _classify_figure_type(self, shape, ocr_text: str) -> str:
        """启发式判断图片类型"""
        try:
            # 基于尺寸比例判断
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                ratio = shape.width / shape.height if shape.height > 0 else 1
                if 1.2 <= ratio <= 2.5:  # 常见图表比例
                    return "chart"
            
            # 基于OCR文本内容判断
            if ocr_text:
                chart_keywords = ["图", "表", "数据", "统计", "分析", "趋势", "%", "比例"]
                if any(keyword in ocr_text for keyword in chart_keywords):
                    return "chart"
                
                diagram_keywords = ["流程", "架构", "结构", "示意", "框图"]
                if any(keyword in ocr_text for keyword in diagram_keywords):
                    return "diagram"
            
            return "photo"
            
        except:
            return "unknown"
